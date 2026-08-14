"""
EvolvixOS — Document Processor Skill
Read and write PDFs, Word docs, Excel, PowerPoint, text files.
100% local. Zero tokens. Zero cloud.

Pip: pip install PyPDF2 python-docx openpyxl python-pptx pdfplumber
License: BSD (PyPDF2), MIT (python-docx), MIT (openpyxl), MIT (python-pptx)
"""

import os
import json
import time
from pathlib import Path
from typing import Optional, List
from rich.console import Console

console = Console()


class Skill:
    """Document processor — PDF, Word, Excel, PowerPoint. Free, local."""

    def __init__(self, config: dict = None):
        self.config = config or {}
        self.output_dir = Path(self.config.get("output_dir", "./output/documents"))
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def run(self, args: dict) -> str:
        action = args.get("action", "read")

        if action == "read":
            return self.read_document(args.get("file", ""))
        elif action == "read_pdf":
            return self.read_pdf(args.get("file", ""))
        elif action == "read_docx":
            return self.read_docx(args.get("file", ""))
        elif action == "read_excel":
            return self.read_excel(args.get("file", ""))
        elif action == "read_pptx":
            return self.read_pptx(args.get("file", ""))
        elif action == "write_pdf":
            return self.write_pdf(args.get("text", ""), args.get("filename", ""))
        elif action == "write_docx":
            return self.write_docx(args.get("text", ""), args.get("filename", ""))
        elif action == "write_excel":
            return self.write_excel(args.get("data", []), args.get("filename", ""))
        elif action == "write_pptx":
            return self.write_pptx(args.get("slides", []), args.get("filename", ""))
        elif action == "convert":
            return self.convert(args.get("file", ""), args.get("to_format", "pdf"))
        elif action == "merge_pdf":
            return self.merge_pdf(args.get("files", []), args.get("output", ""))
        elif action == "split_pdf":
            return self.split_pdf(args.get("file", ""), args.get("pages", []))
        elif action == "extract_images":
            return self.extract_images_from_pdf(args.get("file", ""))
        else:
            return ("Unknown action. Use: read, read_pdf, read_docx, read_excel, "
                    "read_pptx, write_pdf, write_docx, write_excel, write_pptx, "
                    "convert, merge_pdf, split_pdf, extract_images")

    def read_document(self, file_path: str) -> str:
        if not file_path or not os.path.exists(file_path):
            return "Error: File not found."

        ext = file_path.rsplit(".", 1)[-1].lower()

        if ext == "pdf":
            return self.read_pdf(file_path)
        elif ext in ("docx", "doc"):
            return self.read_docx(file_path)
        elif ext in ("xlsx", "xls"):
            return self.read_excel(file_path)
        elif ext in ("pptx", "ppt"):
            return self.read_pptx(file_path)
        elif ext in ("txt", "md", "rst", "csv", "json", "yaml", "yml"):
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()[:50000]
        else:
            return f"Unsupported format: .{ext}"

    def read_pdf(self, file_path: str) -> str:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n\n"
            return text[:50000]
        except ImportError:
            try:
                import pdfplumber
                text = ""
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() + "\n\n"
                return text[:50000]
            except ImportError:
                return "Error: pip install PyPDF2 pdfplumber"
        except Exception as e:
            return f"Error reading PDF: {e}"

    def read_docx(self, file_path: str) -> str:
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs)[:50000]
        except ImportError:
            return "Error: pip install python-docx"
        except Exception as e:
            return f"Error reading docx: {e}"

    def read_excel(self, file_path: str) -> str:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True)
            result = {}
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(max_row=100, values_only=True):
                    rows.append([str(c) if c is not None else "" for c in row])
                result[sheet_name] = rows
            wb.close()
            return json.dumps(result, indent=2)[:10000]
        except ImportError:
            return "Error: pip install openpyxl"
        except Exception as e:
            return f"Error reading Excel: {e}"

    def read_pptx(self, file_path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text.append(shape.text_frame.text)
                slides.append(f"--- Slide {i} ---\n" + "\n".join(text))
            return "\n\n".join(slides)[:50000]
        except ImportError:
            return "Error: pip install python-pptx"
        except Exception as e:
            return f"Error reading PPTX: {e}"

    def write_pdf(self, text: str, filename: str = "") -> str:
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            if not filename:
                filename = f"doc_{int(time.time())}.pdf"
            filepath = self.output_dir / filename
            c = canvas.Canvas(str(filepath), pagesize=letter)
            y = 750
            for line in text.split("\n"):
                c.drawString(50, y, line[:100])
                y -= 15
                if y < 50:
                    c.showPage()
                    y = 750
            c.save()
            return f"PDF written: {filepath}"
        except ImportError:
            return "Error: pip install reportlab"
        except Exception as e:
            return f"Error writing PDF: {e}"

    def write_docx(self, text: str, filename: str = "") -> str:
        try:
            from docx import Document
            if not filename:
                filename = f"doc_{int(time.time())}.docx"
            filepath = self.output_dir / filename
            doc = Document()
            for para in text.split("\n"):
                doc.add_paragraph(para)
            doc.save(str(filepath))
            return f"DOCX written: {filepath}"
        except ImportError:
            return "Error: pip install python-docx"
        except Exception as e:
            return f"Error: {e}"

    def write_excel(self, data: list, filename: str = "") -> str:
        try:
            from openpyxl import Workbook
            if not filename:
                filename = f"data_{int(time.time())}.xlsx"
            filepath = self.output_dir / filename
            wb = Workbook()
            ws = wb.active
            for row in data:
                ws.append(row)
            wb.save(str(filepath))
            return f"Excel written: {filepath}"
        except Exception as e:
            return f"Error: {e}"

    def write_pptx(self, slides: list, filename: str = "") -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches
            if not filename:
                filename = f"slides_{int(time.time())}.pptx"
            filepath = self.output_dir / filename
            prs = Presentation()
            for slide_text in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_text.get("title", "")
                if "body" in slide_text:
                    slide.shapes.placeholders[1].text = slide_text["body"]
            prs.save(str(filepath))
            return f"PPTX written: {filepath}"
        except Exception as e:
            return f"Error: {e}"

    def convert(self, file_path: str, to_format: str = "pdf") -> str:
        text = self.read_document(file_path)
        if text.startswith("Error"):
            return text
        base = Path(file_path).stem
        if to_format == "pdf":
            return self.write_pdf(text, f"{base}.pdf")
        elif to_format == "docx":
            return self.write_docx(text, f"{base}.docx")
        elif to_format == "txt":
            out = self.output_dir / f"{base}.txt"
            out.write_text(text)
            return f"Converted: {file_path} → {out}"
        elif to_format == "json":
            out = self.output_dir / f"{base}.json"
            out.write_text(json.dumps({"content": text}, indent=2))
            return f"Converted: {file_path} → {out}"
        else:
            return f"Unsupported format: {to_format}"

    def merge_pdf(self, files: list, output: str = "") -> str:
        try:
            from PyPDF2 import PdfWriter, PdfReader
            if not output:
                output = str(self.output_dir / f"merged_{int(time.time())}.pdf")
            writer = PdfWriter()
            for f in files:
                if os.path.exists(f):
                    reader = PdfReader(f)
                    for page in reader.pages:
                        writer.add_page(page)
            with open(output, "wb") as f:
                writer.write(f)
            return f"Merged {len(files)} PDFs → {output}"
        except Exception as e:
            return f"Error: {e}"

    def split_pdf(self, file_path: str, pages: list = None) -> str:
        try:
            from PyPDF2 import PdfWriter, PdfReader
            reader = PdfReader(file_path)
            results = []
            for i, page in enumerate(reader.pages):
                writer = PdfWriter()
                writer.add_page(page)
                out = self.output_dir / f"{Path(file_path).stem}_page_{i+1}.pdf"
                with open(out, "wb") as f:
                    writer.write(f)
                results.append(str(out))
            return f"Split {file_path} into {len(results)} pages"
        except Exception as e:
            return f"Error: {e}"

    def extract_images_from_pdf(self, file_path: str) -> str:
        try:
            from PyPDF2 import PdfReader
            import io
            from PIL import Image
            reader = PdfReader(file_path)
            count = 0
            for page in reader.pages:
                for img in page.images:
                    img_path = self.output_dir / f"img_{count}_{int(time.time())}.png"
                    with open(img_path, "wb") as f:
                        f.write(img.data)
                    count += 1
            return f"Extracted {count} images from {file_path}"
        except Exception as e:
            return f"Error: {e}"
